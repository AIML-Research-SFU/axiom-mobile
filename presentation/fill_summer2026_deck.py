#!/usr/bin/env python3
"""Fill in AXIOM-Mobile-Summer-2026.pptx with real content from this
session's work. Only touches placeholders that were empty/generic
("Double-click to edit", "N%", stray template text) -- slides that
already had real, deliberate content (title, team slides, agenda,
thank-you) are left untouched."""
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn


def replace_picture(placeholder, image_path):
    """Swap the image of an already-populated picture placeholder --
    PlaceholderPicture (already has an image) has no .insert_picture(),
    unlike an empty PicturePlaceholder, so the image part + blip r:embed
    relationship have to be swapped directly."""
    image_part, rId = placeholder.part.get_or_add_image_part(image_path)
    blip = placeholder._element.blipFill.find(qn("a:blip"))
    blip.set(qn("r:embed"), rId)

SRC = "/Users/arieltyson/Downloads/AXIOM-Mobile-Summer-2026.pptx"
DST = "/Users/arieltyson/Downloads/AXIOM-Mobile-Summer-2026-filled.pptx"
SCREENSHOT = "/tmp/axiom_demo_screenshot2.png"

prs = Presentation(SRC)
slides = list(prs.slides)


def set_placeholder_text(slide, idx, text, keep_size=True):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    # Preserve the run-level font size of the first run if present
    size = None
    if keep_size and tf.paragraphs and tf.paragraphs[0].runs:
        size = tf.paragraphs[0].runs[0].font.size
    tf.clear()
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    run = p0.add_run()
    run.text = lines[0]
    if size:
        run.font.size = size
    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if size:
            r.font.size = size


def set_bullets(slide, idx, bullets, keep_size=True):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    size = None
    if keep_size and tf.paragraphs and tf.paragraphs[0].runs:
        size = tf.paragraphs[0].runs[0].font.size
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = b
        if size:
            run.font.size = size


# ── Slide 6 (index 5): Quote ────────────────────────────────────────────
set_placeholder_text(
    slides[5], 1,
    "What is the minimum labeled training data a model needs to reason "
    "accurately about what’s on your screen—entirely on-device?"
)

# ── Slide 7 (index 6): What is AXIOM-Mobile? ────────────────────────────
set_placeholder_text(
    slides[6], 1,
    "AXIOM-Mobile is a fully on-device visual question-answering system "
    "for iOS. Show it a screenshot, ask a question in plain language, "
    "and a Core ML model answers instantly—no network call, no cloud, "
    "sub-15ms inference on a real iPhone.\n\n"
    "This summer’s question: how little labeled data does a model "
    "actually need to do this well?"
)
pic_ph = slides[6].placeholders[21]
replace_picture(pic_ph, SCREENSHOT)

# ── Slide 8 (index 7): Topic divider ────────────────────────────────────
set_placeholder_text(slides[7], 1, "Where we left off")

# ── Slide 9 (index 8): Bullet + Graphic ─────────────────────────────────
set_placeholder_text(slides[8], 0, "Where Semester 2 Left Off")
set_bullets(slides[8], 1, [
    "Dataset v3: 751 examples, 93% with zero manual labeling",
    "A pretrained-backbone + LoRA model called “a wash, not a win”—on a single run, no seed variance",
    "A 240-run selection-strategy sweep—but only ever against a memorization heuristic, not a real trainable model",
    "The real trainable-model sweep: attempted, abandoned at 50/240 runs, blamed on the laptop overheating",
    "Energy and memory: never measured, any model, any semester",
])

# ── Slide 10 (index 9): Topic divider ───────────────────────────────────
set_placeholder_text(slides[9], 1, "Where we are now")

# ── Slide 11 (index 10): Bullet + Graphic ───────────────────────────────
set_placeholder_text(slides[10], 0, "What We Actually Found")
set_bullets(slides[10], 1, [
    "The “overheating” was a bug: training never used the GPU. Fixed—the real sweep now runs in under 30 minutes",
    "The “wash, not a win” rested on n=1. At 5 seeds, both pretrained models significantly beat from-scratch (+4.6 to +5.1pp test EM)",
    "Fixed the LoRA text-tower path v4 gave up on—a fixable tracing issue, not a dead end",
    "Real sweep result: Uncertainty/Diversity selection collapse to ~0% below budget 600—invisible on a heuristic; KG-guided wins at practical budgets",
    "Dataset grew to 797 examples via genuinely new content (Safari, Contacts) after finding a stale assumption about simulator apps",
    "Measured energy & memory for the first time, ever—memory passes comfortably, energy is a real narrow split (4.8–5.2%/hr)",
])

# ── Slide 12 (index 11): Topic divider ──────────────────────────────────
set_placeholder_text(slides[11], 1, "What comes next")

# ── Slide 13 (index 12): Bullet + Graphic ───────────────────────────────
set_placeholder_text(slides[12], 0, "Next Steps")
set_bullets(slides[12], 1, [
    "Re-run the KG-guided sweep against the freshly rebuilt v4 knowledge graph",
    "Re-run the architecture comparison on the full committed split, manual examples included",
    "On-device integration for the text-LoRA model (Swift tokenizer)",
    "Investigate why pretrained backbones help less than hoped—try a UI/document-pretrained backbone instead of ImageNet",
    "Close the quality gap toward the 70% target—still the open question",
])

# ── Slide 14 (index 13): Large Number ───────────────────────────────────
set_placeholder_text(slides[13], 22, "32.0%")
set_placeholder_text(
    slides[13], 21,
    "tiny_multimodal_v1 test EM — 106 KB, trained from scratch, 5-seed mean"
)

# ── Slide 15 (index 14): remove stray "Apple Confidential" text ────────
slide15 = slides[14]
for shape in list(slide15.shapes):
    if shape.has_text_frame and "Apple Confidential" in shape.text_frame.text:
        shape._element.getparent().remove(shape._element)

prs.save(DST)
print(f"Saved: {DST}")
